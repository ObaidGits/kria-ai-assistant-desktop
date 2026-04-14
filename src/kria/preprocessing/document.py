"""
Document & Office Preprocessing Module
========================================
Parse PDF, DOCX, XLSX, PPTX, and CSV via Microsoft MarkItDown (primary)
with fallback to PyMuPDF / python-docx / pandas.  Extracts clean markdown
text and structural metadata — **no page-image generation**.
"""
from __future__ import annotations

import asyncio
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger("kria.preprocessing.document")

_MARKITDOWN_AVAILABLE: Optional[bool] = None


def _check_markitdown() -> bool:
    global _MARKITDOWN_AVAILABLE
    if _MARKITDOWN_AVAILABLE is None:
        try:
            import markitdown  # noqa: F401
            _MARKITDOWN_AVAILABLE = True
        except ImportError:
            _MARKITDOWN_AVAILABLE = False
            logger.info("markitdown not installed — falling back to PyMuPDF/python-docx/pandas")
    return _MARKITDOWN_AVAILABLE


async def preprocess_document(
    source: str,
    *,
    content: Optional[bytes] = None,
    max_tokens: int = 3500,
) -> "PreprocessedPayload":
    """Convert a document file to token-budgeted markdown text.

    Supports: PDF, DOCX, XLSX, PPTX, CSV.
    Primary parser: MarkItDown → clean markdown.
    Fallback: PyMuPDF (PDF), python-docx (DOCX), pandas (XLSX/CSV).
    """
    from kria.preprocessing.dispatcher import PreprocessedPayload
    from kria.preprocessing.token_budget import estimate_tokens, smart_crop

    ext = Path(source).suffix.lower()

    def _extract() -> tuple[str, dict]:
        """Extract text + metadata (runs in thread)."""
        # Write content to temp file if provided as bytes (MarkItDown needs a path)
        file_path = source
        tmp_path = None
        if content:
            import tempfile
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
            tmp.write(content)
            tmp.flush()
            tmp.close()
            file_path = tmp.name
            tmp_path = tmp.name

        try:
            text, meta = _extract_with_path(file_path, ext)
        finally:
            if tmp_path:
                try:
                    Path(tmp_path).unlink()
                except OSError:
                    pass
        return text, meta

    text, meta = await asyncio.to_thread(_extract)
    meta["path"] = source

    text, truncated = smart_crop(text, max_tokens)
    tokens = estimate_tokens(text)

    return PreprocessedPayload(
        text=text,
        token_estimate=tokens,
        source_type="document",
        metadata=meta,
        truncated=truncated,
    )


def _extract_with_path(file_path: str, ext: str) -> tuple[str, dict]:
    """Dispatch extraction by extension, MarkItDown first."""
    if _check_markitdown():
        try:
            return _extract_markitdown(file_path, ext)
        except Exception as exc:
            logger.warning("MarkItDown failed for %s, falling back: %s", file_path, exc)

    # Fallbacks per extension
    if ext == ".pdf":
        return _extract_pdf_fallback(file_path)
    if ext in {".docx", ".doc"}:
        return _extract_docx_fallback(file_path)
    if ext in {".xlsx", ".xls"}:
        return _extract_xlsx_fallback(file_path)
    if ext == ".csv":
        return _extract_csv_fallback(file_path)
    if ext == ".pptx":
        return _extract_pptx_fallback(file_path)

    # Last resort: read as text
    try:
        text = Path(file_path).read_text(encoding="utf-8", errors="replace")
        return text, {"parser": "raw_text"}
    except Exception:
        return "", {"parser": "none", "error": "unsupported format"}


# ---------------------------------------------------------------------------
# MarkItDown (primary)
# ---------------------------------------------------------------------------

def _extract_markitdown(file_path: str, ext: str) -> tuple[str, dict]:
    from markitdown import MarkItDown

    md = MarkItDown()
    result = md.convert(file_path)
    text = result.text_content or ""

    # Build structural metadata
    import re
    headings = re.findall(r"^#{1,6}\s+(.+)$", text, re.MULTILINE)
    table_count = text.count("|---") + text.count("| ---")

    meta = {
        "parser": "markitdown",
        "format": ext.lstrip("."),
        "heading_count": len(headings),
        "table_count": table_count,
        "char_count": len(text),
    }
    if hasattr(result, "title") and result.title:
        meta["title"] = result.title

    return text, meta


# ---------------------------------------------------------------------------
# Fallback parsers
# ---------------------------------------------------------------------------

def _extract_pdf_fallback(file_path: str) -> tuple[str, dict]:
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    pages = []
    for page in doc:
        t = page.get_text().strip()
        if t:
            pages.append(t)
    total = len(doc)
    doc.close()
    text = "\n\n".join(pages)
    return text, {
        "parser": "pymupdf",
        "format": "pdf",
        "total_pages": total,
        "extracted_pages": len(pages),
        "char_count": len(text),
    }


def _extract_docx_fallback(file_path: str) -> tuple[str, dict]:
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Tables as markdown
    table_parts: list[str] = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            # Add header separator after first row
            header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
            rows.insert(1, header_sep)
            table_parts.append("\n".join(rows))

    text = "\n\n".join(paragraphs)
    if table_parts:
        text += "\n\n" + "\n\n".join(table_parts)

    return text, {
        "parser": "python-docx",
        "format": "docx",
        "paragraph_count": len(paragraphs),
        "table_count": len(table_parts),
        "char_count": len(text),
    }


def _extract_xlsx_fallback(file_path: str) -> tuple[str, dict]:
    import pandas as pd

    xls = pd.ExcelFile(file_path)
    parts: list[str] = []
    total_rows = 0

    for sheet_name in xls.sheet_names[:5]:  # Cap at 5 sheets
        df = pd.read_excel(xls, sheet_name=sheet_name, nrows=200)
        total_rows += len(df)
        md_table = df.to_markdown(index=False)
        parts.append(f"## Sheet: {sheet_name}\n\n{md_table}")

    text = "\n\n".join(parts)
    return text, {
        "parser": "pandas",
        "format": "xlsx",
        "sheet_count": len(xls.sheet_names),
        "total_rows": total_rows,
        "char_count": len(text),
    }


def _extract_csv_fallback(file_path: str) -> tuple[str, dict]:
    import pandas as pd

    df = pd.read_csv(file_path, nrows=200)
    text = df.to_markdown(index=False)
    return text, {
        "parser": "pandas",
        "format": "csv",
        "rows": len(df),
        "columns": list(df.columns),
        "char_count": len(text),
    }


def _extract_pptx_fallback(file_path: str) -> tuple[str, dict]:
    """Basic PPTX text extraction without MarkItDown."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            if texts:
                slides.append(f"## Slide {i}\n\n" + "\n".join(texts))

        text = "\n\n".join(slides)
        return text, {
            "parser": "python-pptx",
            "format": "pptx",
            "slide_count": len(prs.slides),
            "char_count": len(text),
        }
    except ImportError:
        return "", {"parser": "none", "format": "pptx", "error": "python-pptx not installed"}
